"""
Document processor for extracting text and creating chunks.

PDF extraction uses pymupdf (fitz) for better layout handling.
PPTX extraction uses python-pptx (slide-level chunking preserves context).
PyPDF2 kept as fallback only.
"""

import os
import re
import logging
from pathlib import Path
from typing import List, Dict, Any, Optional

import docx
from bs4 import BeautifulSoup

# pymupdf is the primary PDF extractor (already in requirements)
try:
    import fitz as _fitz  # pymupdf
    _FITZ_AVAILABLE = True
except ImportError:
    import PyPDF2 as _PyPDF2  # fallback
    _FITZ_AVAILABLE = False

# PPTX support (python-pptx)
try:
    from pptx import Presentation as _Presentation
    _PPTX_AVAILABLE = True
except ImportError:
    _PPTX_AVAILABLE = False

from extensions import ProgressTracker

logger = logging.getLogger(__name__)

class DocumentProcessor:
    """Enhanced document processor that supports both normal and graph RAG modes."""

    def __init__(self, config):
        self.chunk_size = config["chunk_size"]
        self.chunk_overlap = config["chunk_overlap"]
        self.rag_mode = config.get("rag_mode", "normal")

        self.supported_extensions = {'.pdf', '.txt', '.docx', '.html', '.htm'}
        if _PPTX_AVAILABLE:
            self.supported_extensions.add('.pptx')
        
        logger.info(f"DocumentProcessor initialized with chunk_size={self.chunk_size}, rag_mode={self.rag_mode}")
    
    def is_supported_file(self, file_path: str) -> bool:
        """Check if the file type is supported."""
        path = Path(file_path)
        return path.suffix.lower() in self.supported_extensions
    
    def process_document(self, file_path: str, 
                        progress_tracker: Optional[ProgressTracker] = None) -> List[Dict[str, Any]]:
        """
        Complete document processing pipeline: extract text and create chunks.
        
        This is a convenience method that combines extract_text_from_file() and chunk_text().
        
        Args:
            file_path: Path to the document file
            progress_tracker: Optional progress tracker for UI updates
            
        Returns:
            List of chunk dictionaries with 'text' and 'metadata' fields
        """
        if progress_tracker:
            progress_tracker.update(0, 100, status="extracting", 
                                   message="Extracting text from document")
        
        # Step 1: Extract text from file
        text = self.extract_text_from_file(file_path, progress_tracker)
        
        if not text or not text.strip():
            logger.warning(f"No text extracted from {file_path}")
            return []
        
        if progress_tracker:
            progress_tracker.update(50, 100, status="chunking", 
                                   message="Creating text chunks")
        
        # Step 2: Create metadata for the document
        filename = Path(file_path).name
        metadata = {
            "source": filename,
            "file_path": file_path,
            "file_type": Path(file_path).suffix.lower()
        }
        
        # Step 3: Chunk the text
        chunks = self.chunk_text(text, metadata, progress_tracker)
        
        if progress_tracker:
            progress_tracker.update(100, 100, status="complete", 
                                   message=f"Created {len(chunks)} chunks")
        
        logger.info(f"Processed {filename}: extracted {len(text)} chars, created {len(chunks)} chunks")
        return chunks
    
    def extract_text_from_file(self, file_path: str, 
                             progress_tracker: Optional[ProgressTracker] = None) -> str:
        """Extract text from various file formats with comprehensive error handling."""
        path = Path(file_path)
        
        if not path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
        
        extension = path.suffix.lower()
        
        if extension not in self.supported_extensions:
            raise ValueError(f"Unsupported file format: {extension}. Supported formats: {self.supported_extensions}")
        
        if progress_tracker:
            progress_tracker.update(0, 100, status="extracting", 
                                   message=f"Extracting text from {path.name}",
                                   current_file=path.name)
        
        try:
            if extension == '.pdf':
                return self._extract_from_pdf(file_path, progress_tracker)
            elif extension == '.txt':
                return self._extract_from_txt(file_path, progress_tracker)
            elif extension == '.docx':
                return self._extract_from_docx(file_path, progress_tracker)
            elif extension in ['.html', '.htm']:
                return self._extract_from_html(file_path, progress_tracker)
            elif extension == '.pptx':
                return self._extract_from_pptx(file_path, progress_tracker)
            else:
                raise ValueError(f"Unsupported file format: {extension}")
        
        except Exception as e:
            logger.error(f"Error extracting text from {file_path}: {str(e)}")
            if progress_tracker:
                progress_tracker.update(100, 100, status="error", 
                                       message=f"Error extracting text: {str(e)}")
            raise
    
    def _extract_from_pdf(self, file_path: str,
                         progress_tracker: Optional[ProgressTracker] = None) -> str:
        """Extract text from PDF using pymupdf (better layout/table handling than PyPDF2)."""
        if _FITZ_AVAILABLE:
            return self._extract_pdf_fitz(file_path, progress_tracker)
        return self._extract_pdf_pypdf2_fallback(file_path, progress_tracker)

    def _extract_pdf_fitz(self, file_path: str,
                          progress_tracker: Optional[ProgressTracker] = None) -> str:
        """pymupdf extraction — preserves column order, handles tables as text blocks."""
        try:
            doc = _fitz.open(file_path)
            total_pages = len(doc)
            if total_pages == 0:
                return ""

            parts = []
            for i, page in enumerate(doc):
                # "text" mode with "blocks" sorts by reading order
                page_text = page.get_text("text")
                if page_text.strip():
                    parts.append(page_text)
                if progress_tracker:
                    pct = int(((i + 1) / total_pages) * 100)
                    progress_tracker.update(pct, 100, message=f"Extracting page {i + 1}/{total_pages}")
            doc.close()

            extracted = "\n".join(parts)
            if not extracted.strip():
                logger.warning(f"No text extracted from PDF: {file_path}")
                return ""
            logger.info(f"pymupdf extracted {len(extracted)} chars from {total_pages}-page PDF")
            return extracted
        except Exception as e:
            logger.error(f"pymupdf extraction failed for {file_path}: {e}")
            raise RuntimeError(f"Failed to extract PDF text: {e}")

    def _extract_pdf_pypdf2_fallback(self, file_path: str,
                                     progress_tracker: Optional[ProgressTracker] = None) -> str:
        """PyPDF2 fallback when pymupdf is not available."""
        try:
            import PyPDF2
            text_parts = []
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                total_pages = len(reader.pages)
                for i, page in enumerate(reader.pages):
                    try:
                        t = page.extract_text()
                        if t and t.strip():
                            text_parts.append(t)
                    except Exception:
                        pass
                    if progress_tracker:
                        pct = int(((i + 1) / total_pages) * 100)
                        progress_tracker.update(pct, 100, message=f"Extracting page {i + 1}/{total_pages}")
            return "\n".join(text_parts)
        except Exception as e:
            raise RuntimeError(f"Failed to extract PDF text: {e}")

    def _extract_from_pptx(self, file_path: str,
                            progress_tracker: Optional[ProgressTracker] = None) -> str:
        """
        Extract text from PPTX slide-by-slide.

        Each slide becomes one logical block: title + body + table cells.
        This preserves the atomic context of each slide (e.g. 'Market Size' slide
        keeps its heading, data points, and footer together).
        """
        if not _PPTX_AVAILABLE:
            raise RuntimeError(
                "python-pptx is not installed. Run: pip install python-pptx"
            )
        try:
            prs = _Presentation(file_path)
            slide_texts = []
            total = len(prs.slides)

            for i, slide in enumerate(prs.slides):
                slide_parts = []

                # Slide title (if present)
                if slide.shapes.title and slide.shapes.title.text.strip():
                    slide_parts.append(f"## {slide.shapes.title.text.strip()}")

                for shape in slide.shapes:
                    # Text frames (body text)
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = para.text.strip()
                            if t and t != slide_parts[0].lstrip("# ") if slide_parts else True:
                                slide_parts.append(t)

                    # Tables — serialized as markdown-style rows
                    if shape.has_table:
                        for row in shape.table.rows:
                            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                            if cells:
                                slide_parts.append(" | ".join(cells))

                if slide_parts:
                    slide_texts.append(f"[Slide {i + 1}]\n" + "\n".join(slide_parts))

                if progress_tracker:
                    pct = int(((i + 1) / total) * 100)
                    progress_tracker.update(pct, 100, message=f"Extracting slide {i + 1}/{total}")

            extracted = "\n\n".join(slide_texts)
            logger.info(f"python-pptx extracted {len(extracted)} chars from {total}-slide PPTX")
            return extracted
        except Exception as e:
            logger.error(f"PPTX extraction failed for {file_path}: {e}")
            raise RuntimeError(f"Failed to extract PPTX text: {e}")
    
    def _extract_from_txt(self, file_path: str, 
                         progress_tracker: Optional[ProgressTracker] = None) -> str:
        """Extract text from TXT file with encoding detection."""
        try:
            # Try different encodings
            encodings = ['utf-8', 'utf-16', 'latin-1', 'cp1252']
            
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding, errors='replace') as file:
                        text = file.read()
                        
                        if progress_tracker:
                            progress_tracker.update(100, 100, status="complete", 
                                                   message="Text file extraction complete")
                        
                        logger.info(f"Successfully extracted {len(text)} characters from TXT file using {encoding} encoding")
                        return text
                        
                except UnicodeDecodeError:
                    continue
            
            # If all encodings fail, try with errors='ignore'
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                text = file.read()
                logger.warning(f"Used UTF-8 with error ignoring for file: {file_path}")
                return text
                
        except Exception as e:
            logger.error(f"Error reading TXT file {file_path}: {str(e)}")
            raise RuntimeError(f"Failed to extract text from TXT file: {str(e)}")
    
    def _extract_from_docx(self, file_path: str, 
                          progress_tracker: Optional[ProgressTracker] = None) -> str:
        """Extract text from DOCX file with comprehensive content extraction."""
        try:
            doc = docx.Document(file_path)
            text_parts = []
            
            # Extract from paragraphs
            paragraphs = doc.paragraphs
            total_paragraphs = len(paragraphs)
            
            logger.debug(f"Extracting text from {total_paragraphs} paragraphs")
            
            for i, paragraph in enumerate(paragraphs):
                if paragraph.text.strip():  # Only add non-empty paragraphs
                    text_parts.append(paragraph.text.strip())
                
                if progress_tracker and total_paragraphs > 0:
                    progress_percentage = int(((i + 1) / total_paragraphs) * 50)  # 50% for paragraphs
                    progress_tracker.update(progress_percentage, 100, 
                                           message=f"Extracting paragraph {i + 1} of {total_paragraphs}")
            
            # Extract from tables
            tables = doc.tables
            if tables:
                logger.debug(f"Extracting text from {len(tables)} tables")
                
                for table_idx, table in enumerate(tables):
                    table_text = []
                    for row in table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_text.append(cell.text.strip())
                        if row_text:
                            table_text.append(" | ".join(row_text))
                    
                    if table_text:
                        text_parts.append("\n".join(table_text))
                    
                    if progress_tracker:
                        progress_percentage = 50 + int(((table_idx + 1) / len(tables)) * 50)  # 50% for tables
                        progress_tracker.update(progress_percentage, 100, 
                                               message=f"Extracting table {table_idx + 1} of {len(tables)}")
            
            extracted_text = "\n\n".join(text_parts)
            
            if not extracted_text.strip():
                logger.warning(f"No text extracted from DOCX: {file_path}")
                return ""
            
            logger.info(f"Successfully extracted {len(extracted_text)} characters from DOCX")
            return extracted_text
            
        except Exception as e:
            logger.error(f"Error reading DOCX file {file_path}: {str(e)}")
            raise RuntimeError(f"Failed to extract text from DOCX file: {str(e)}")
    
    def _extract_from_html(self, file_path: str, 
                          progress_tracker: Optional[ProgressTracker] = None) -> str:
        """Extract text from HTML file with improved content extraction."""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='replace') as file:
                html_content = file.read()
            
            if progress_tracker:
                progress_tracker.update(25, 100, message="Parsing HTML content")
            
            # Parse HTML
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # Remove script and style elements
            for script in soup(["script", "style", "nav", "header", "footer"]):
                script.extract()
            
            if progress_tracker:
                progress_tracker.update(50, 100, message="Extracting text content")
            
            # Get text content
            text = soup.get_text()
            
            if progress_tracker:
                progress_tracker.update(75, 100, message="Cleaning extracted text")
            
            # Clean up the text
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            cleaned_text = '\n'.join(chunk for chunk in chunks if chunk)
            
            if progress_tracker:
                progress_tracker.update(100, 100, message="HTML extraction complete")
            
            if not cleaned_text.strip():
                logger.warning(f"No text extracted from HTML: {file_path}")
                return ""
            
            logger.info(f"Successfully extracted {len(cleaned_text)} characters from HTML")
            return cleaned_text
            
        except Exception as e:
            logger.error(f"Error reading HTML file {file_path}: {str(e)}")
            raise RuntimeError(f"Failed to extract text from HTML file: {str(e)}")
    
    # ── heading patterns: markdown, ALL CAPS labels, numbered sections, slide headers ──
    _HEADING_RE = re.compile(
        r'^(?:'
        r'#{1,4}\s+.+|'                          # ## Markdown heading
        r'\[Slide \d+\]|'                         # [Slide N] from PPTX extractor
        r'[A-Z][A-Z0-9 \t\-:]{4,}$|'            # ALL CAPS LABEL
        r'\d{1,2}\.(?:\d+\.?)?\s+[A-Z].+|'      # 1. or 2.1 Numbered section
        r'(?:SECTION|CHAPTER|PART)\s+\w+'        # SECTION X / CHAPTER Y
        r')',
        re.MULTILINE,
    )

    def chunk_text(self, text: str, metadata: dict,
                   progress_tracker: Optional[ProgressTracker] = None) -> List[Dict[str, Any]]:
        """
        Structure-aware chunking.

        1. Detect section headings (markdown, ALL CAPS, numbered, slide markers).
        2. Group text into sections; each section header is stored as section_path.
        3. Sub-split sections larger than chunk_size using sentence boundaries.
        4. Prepend the heading path as a contextual header on the chunk text so the
           embedding carries section context (the 'late chunking' analogue without
           changing the embedding model).

        Falls back to sentence-based splitting when no headings are detected.
        """
        if not text or not text.strip():
            logger.warning("Empty text provided for chunking")
            return []

        if progress_tracker:
            progress_tracker.update(0, 100, status="chunking", message="Splitting text into chunks")

        effective_chunk_size = self.chunk_size
        if self.rag_mode == "graph":
            effective_chunk_size = int(self.chunk_size * 1.5)

        # Detect sections on line-normalized (but not whitespace-collapsed) text
        # so heading patterns can match against line starts.
        line_normalized = text.replace('\r\n', '\n').replace('\r', '\n')
        sections = self._split_into_sections(line_normalized)

        # If no sections detected, fall back to sentence-based approach on cleaned text
        if len(sections) == 1 and sections[0][0] is None:
            return self._chunk_by_sentences(self._clean_text(text), metadata, effective_chunk_size, progress_tracker)

        chunks: List[Dict[str, Any]] = []
        total_sections = len(sections)

        for sec_idx, (heading, body) in enumerate(sections):
            body = self._clean_text(body)
            if not body.strip():
                continue

            section_path = heading or metadata.get("source", "")
            contextual_header = f"[Section: {section_path}]\n" if section_path else ""

            if len(body) <= effective_chunk_size:
                # Whole section fits in one chunk
                chunk_text = contextual_header + body
                chunk_meta = metadata.copy()
                chunk_meta.update({
                    "chunk_index": len(chunks),
                    "section_path": section_path,
                    "contextual_header": contextual_header.strip(),
                    "rag_mode": self.rag_mode,
                    "chunk_size": len(chunk_text),
                    "chunk_text_preview": body[:100],
                })
                chunks.append({"text": chunk_text, "metadata": chunk_meta})
            else:
                # Sub-split large section by sentences
                sub_chunks = self._chunk_by_sentences(body, metadata, effective_chunk_size)
                for sub in sub_chunks:
                    sub["text"] = contextual_header + sub["text"]
                    sub["metadata"]["section_path"] = section_path
                    sub["metadata"]["contextual_header"] = contextual_header.strip()
                    sub["metadata"]["chunk_index"] = len(chunks)
                    chunks.append(sub)

            if progress_tracker:
                pct = int(((sec_idx + 1) / total_sections) * 100)
                progress_tracker.update(pct, 100, message=f"Chunking section {sec_idx + 1}/{total_sections}")

        if progress_tracker:
            progress_tracker.update(100, 100, status="chunking_complete",
                                    message=f"Created {len(chunks)} chunks")

        avg = sum(len(c["text"]) for c in chunks) // len(chunks) if chunks else 0
        logger.info(f"Structure-aware chunking: {len(chunks)} chunks, avg {avg} chars (mode: {self.rag_mode})")
        return chunks

    def _split_into_sections(self, text: str) -> List[tuple]:
        """
        Split text on detected headings.

        Returns list of (heading_str | None, body_str) tuples.
        If no headings found, returns [(None, full_text)].
        """
        lines = text.split("\n")
        sections = []
        current_heading = None
        current_body: List[str] = []

        for line in lines:
            stripped = line.strip()
            if stripped and self._HEADING_RE.match(stripped):
                # Save accumulated body under previous heading
                if current_body:
                    sections.append((current_heading, "\n".join(current_body).strip()))
                current_heading = stripped.lstrip("#").strip()
                current_body = []
            else:
                current_body.append(line)

        # Flush last section
        if current_body:
            sections.append((current_heading, "\n".join(current_body).strip()))

        # If only one section with no heading, headings weren't found
        if len(sections) == 1 and sections[0][0] is None:
            return sections

        return [s for s in sections if s[1].strip()]

    def _chunk_by_sentences(self, text: str, metadata: dict,
                             chunk_size: int = None,
                             progress_tracker: Optional[ProgressTracker] = None) -> List[Dict[str, Any]]:
        """Original sentence-boundary chunking — used as sub-splitter for large sections."""
        if chunk_size is None:
            chunk_size = self.chunk_size

        sentences = self._split_into_sentences(text)
        if not sentences:
            return []

        chunks: List[Dict[str, Any]] = []
        current = ""
        sent_count = 0

        for i, sentence in enumerate(sentences):
            sentence = sentence.strip()
            if not sentence:
                continue

            candidate = (current + " " + sentence).strip() if current else sentence

            if len(candidate) > chunk_size and current:
                chunk_meta = metadata.copy()
                chunk_meta.update({
                    "chunk_index": len(chunks),
                    "chunk_text_preview": current[:100],
                    "rag_mode": self.rag_mode,
                    "sentence_count": sent_count,
                    "chunk_size": len(current),
                })
                chunks.append({"text": current, "metadata": chunk_meta})

                overlap = self._get_overlap_sentences(sentences, i, sent_count)
                current = (" ".join(overlap) + " " + sentence).strip() if overlap else sentence
                sent_count = len(overlap) + 1
            else:
                current = candidate
                sent_count += 1

            if progress_tracker:
                pct = int(((i + 1) / len(sentences)) * 100)
                progress_tracker.update(pct, 100, message=f"Sentence {i + 1}/{len(sentences)}")

        if current.strip():
            chunk_meta = metadata.copy()
            chunk_meta.update({
                "chunk_index": len(chunks),
                "chunk_text_preview": current[:100],
                "rag_mode": self.rag_mode,
                "sentence_count": sent_count,
                "chunk_size": len(current),
            })
            chunks.append({"text": current, "metadata": chunk_meta})

        return chunks
    
    def _clean_text(self, text: str) -> str:
        """Clean and normalize text for better processing."""
        # Remove excessive whitespace
        cleaned = re.sub(r'\s+', ' ', text)
        
        # Remove or replace problematic characters
        cleaned = cleaned.replace('\x00', '')  # Remove null bytes
        cleaned = cleaned.replace('\ufeff', '')  # Remove BOM
        
        # Normalize line breaks
        cleaned = cleaned.replace('\r\n', '\n').replace('\r', '\n')
        
        # Remove excessive line breaks
        cleaned = re.sub(r'\n\s*\n\s*\n+', '\n\n', cleaned)
        
        return cleaned.strip()
    
    def _split_into_sentences(self, text: str) -> List[str]:
        """Split text into sentences using improved regex patterns."""
        # Pattern for sentence boundaries
        sentence_pattern = r'(?<=[.!?])\s+(?=[A-Z])|(?<=[.!?])\s*\n+\s*(?=[A-Z])'
        
        # Split by sentence boundaries
        sentences = re.split(sentence_pattern, text)
        
        # Clean and filter sentences
        cleaned_sentences = []
        for sentence in sentences:
            sentence = sentence.strip()
            if sentence and len(sentence) > 10:  # Filter out very short fragments
                cleaned_sentences.append(sentence)
        
        # If regex splitting didn't work well, fall back to simple splitting
        if len(cleaned_sentences) < 2:
            # Simple fallback: split by periods followed by space and capital letter
            sentences = re.split(r'\.(?=\s+[A-Z])', text)
            cleaned_sentences = [s.strip() + '.' for s in sentences if s.strip()]
        
        return cleaned_sentences
    
    def _get_overlap_sentences(self, sentences: List[str], current_index: int, 
                              current_sentence_count: int) -> List[str]:
        """Get sentences for overlap between chunks."""
        if not sentences or current_sentence_count == 0:
            return []
        
        # Calculate number of sentences for overlap based on chunk_overlap setting
        # Estimate average sentence length and calculate overlap sentences
        avg_sentence_length = 100  # Rough estimate
        overlap_sentences_count = max(1, self.chunk_overlap // avg_sentence_length)
        overlap_sentences_count = min(overlap_sentences_count, current_sentence_count - 1)
        
        if overlap_sentences_count <= 0:
            return []
        
        # Get the last N sentences from the current chunk
        start_index = max(0, current_index - overlap_sentences_count)
        return sentences[start_index:current_index]
    
    def validate_chunks(self, chunks: List[Dict[str, Any]]) -> Dict[str, Any]:
        """Validate chunks and return statistics."""
        if not chunks:
            return {
                "valid": False,
                "error": "No chunks provided",
                "total_chunks": 0
            }
        
        stats = {
            "valid": True,
            "total_chunks": len(chunks),
            "total_characters": 0,
            "average_chunk_size": 0,
            "min_chunk_size": float('inf'),
            "max_chunk_size": 0,
            "empty_chunks": 0,
            "chunks_too_large": 0,
            "chunks_too_small": 0
        }
        
        chunk_sizes = []
        
        for i, chunk in enumerate(chunks):
            # Validate chunk structure
            if not isinstance(chunk, dict):
                stats["valid"] = False
                stats["error"] = f"Chunk {i} is not a dictionary"
                return stats
            
            if "text" not in chunk:
                stats["valid"] = False
                stats["error"] = f"Chunk {i} missing 'text' field"
                return stats
            
            if "metadata" not in chunk:
                stats["valid"] = False
                stats["error"] = f"Chunk {i} missing 'metadata' field"
                return stats
            
            # Analyze chunk size
            chunk_text = chunk["text"]
            chunk_size = len(chunk_text)
            chunk_sizes.append(chunk_size)
            
            stats["total_characters"] += chunk_size
            stats["min_chunk_size"] = min(stats["min_chunk_size"], chunk_size)
            stats["max_chunk_size"] = max(stats["max_chunk_size"], chunk_size)
            
            # Count issues
            if chunk_size == 0:
                stats["empty_chunks"] += 1
            elif chunk_size < 50:  # Very small chunks
                stats["chunks_too_small"] += 1
            elif chunk_size > self.chunk_size * 2:  # Very large chunks
                stats["chunks_too_large"] += 1
        
        # Calculate averages
        if chunk_sizes:
            stats["average_chunk_size"] = sum(chunk_sizes) / len(chunk_sizes)
            stats["min_chunk_size"] = min(chunk_sizes)
        else:
            stats["min_chunk_size"] = 0
        
        # Add warnings
        warnings = []
        if stats["empty_chunks"] > 0:
            warnings.append(f"{stats['empty_chunks']} empty chunks found")
        if stats["chunks_too_small"] > 0:
            warnings.append(f"{stats['chunks_too_small']} very small chunks found")
        if stats["chunks_too_large"] > 0:
            warnings.append(f"{stats['chunks_too_large']} very large chunks found")
        
        stats["warnings"] = warnings
        
        return stats
    
    def get_supported_formats(self) -> Dict[str, str]:
        """Get information about supported file formats."""
        return {
            ".pdf": "Portable Document Format - extracts text from all pages",
            ".txt": "Plain text files - supports multiple encodings",
            ".docx": "Microsoft Word documents - extracts text and tables",
            ".html": "HTML web pages - extracts text content, removes scripts/styles",
            ".htm": "HTML web pages - same as .html"
        }
    
    def estimate_processing_time(self, file_path: str) -> Dict[str, Any]:
        """Estimate processing time for a file."""
        try:
            path = Path(file_path)
            file_size = path.stat().st_size
            extension = path.suffix.lower()
            
            # Rough estimates based on file type and size
            time_estimates = {
                ".txt": file_size / (1024 * 1024) * 0.1,  # 0.1 seconds per MB
                ".pdf": file_size / (1024 * 1024) * 2.0,   # 2 seconds per MB
                ".docx": file_size / (1024 * 1024) * 1.0,  # 1 second per MB
                ".html": file_size / (1024 * 1024) * 0.5,  # 0.5 seconds per MB
                ".htm": file_size / (1024 * 1024) * 0.5,   # 0.5 seconds per MB
            }
            
            base_time = time_estimates.get(extension, 1.0)
            
            # Add extra time for graph mode
            if self.rag_mode == "graph":
                base_time *= 3  # Graph extraction takes longer
            
            return {
                "estimated_seconds": max(1, int(base_time)),
                "file_size_mb": round(file_size / (1024 * 1024), 2),
                "complexity": "high" if self.rag_mode == "graph" else "medium"
            }
            
        except Exception as e:
            logger.error(f"Error estimating processing time: {e}")
            return {
                "estimated_seconds": 60,  # Default estimate
                "file_size_mb": 0,
                "complexity": "unknown",
                "error": str(e)
            }
    
    def get_processor_stats(self) -> Dict[str, Any]:
        """Get statistics about the document processor."""
        return {
            "rag_mode": self.rag_mode,
            "chunk_size": self.chunk_size,
            "chunk_overlap": self.chunk_overlap,
            "supported_extensions": list(self.supported_extensions),
            "effective_chunk_size": int(self.chunk_size * 1.5) if self.rag_mode == "graph" else self.chunk_size
        }